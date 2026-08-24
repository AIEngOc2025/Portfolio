from pptx import Presentation

def create_portfolio_presentation():
    prs = Presentation()

    # 1. Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Portfolio : AI Engineer"
    slide.placeholders[1].text = "Concevoir, industrialiser et déployer des solutions IA\nPrésenté par Christophe Mpaga"

    # 2. Profil
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Profil : AI Engineer"
    slide.placeholders[1].text = ("• Focus : Industrialisation et passage à l'échelle (POC vers PROD).\n"
                                    "• Approche : Obsession pour la fiabilité, l'observabilité et la conformité.\n"
                                    "• Vision : De l'ingénierie logicielle à l'IA générative.")

    # 3. Piliers Techniques
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Stack Opérationnelle"
    slide.placeholders[1].text = ("• DevOps : CI/CD, Git, GitHub Actions, Docker.\n"
                                    "• MLOps : MLflow, suivi de modèles, évaluation, monitoring.\n"
                                    "• AIOps : Prompt engineering, systèmes multi-agents, red-teaming.\n"
                                    "• CloudOps : GCP (Cloud Run, Vertex AI, GCS), Infrastructure as Code.")

    # 4. Projet Phare : CHSA
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Projet Phare : Agent d'Accueil Hospitalier CHSA"
    slide.placeholders[1].text = ("• Objectif : Réduire la surcharge aux urgences.\n"
                                    "• Démarche : Cadrage métier, architecture modulaire, intégration SIH.\n"
                                    "• MVP & PROD : Fiabilisation, haute disponibilité, interopérabilité (FHIR).")

    # 5. Ingénierie & Pivots
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Ingénierie & Leçons Apprises"
    slide.placeholders[1].text = ("• Pivot Architectural : Échec du design monolithique (timeout 240s) -> Architecture multi-stage modulaire.\n"
                                    "• Modularité : Maintenance découplée et itération indépendante par brique.\n"
                                    "• Futur : Architecture agentique (LangGraph).")

    # 6. Breadth of Projects
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Projets & Applications"
    slide.placeholders[1].text = ("• Fine-tuning LLM (Qwen3) : Industrialisation et serving.\n"
                                    "• Scoring de crédit (LightGBM) : Classification et API conteneurisée.\n"
                                    "• Système RAG (Événements parisiens) : Architecture vecteurs + LLM.")

    # 7. Expérience & Soft Skills
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Expérience & Soft Skills"
    slide.placeholders[1].text = ("• Pédagogie : Enseignement Mathématiques, BDD, ML/IA (vulgarisation).\n"
                                    "• Analyse & Résolution : Audit et refactorisation architecturale.\n"
                                    "• Créativité : Pensée out-of-the-box (agents, modularité).")

    # 8. Éthique & Conformité
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Engagement Éthique et Conformité"
    slide.placeholders[1].text = ("• Protection des données : Anonymisation native et RGPD.\n"
                                    "• Sécurité : Red-teaming et garde-fous contre les hallucinations.\n"
                                    "• Transparence : Auditabilité des systèmes (logs et observabilité).")

    # 9. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion"
    content = slide.placeholders[1]
    content.text = ("• Ingénieur IA orienté production.\n"
                    "• Maîtrise de la chaîne de valeur (de l'idée à la conformité).\n"
                    "• Prêt pour des défis complexes d'industrialisation IA.\n"
                    "• Merci de votre attention. Questions ?")

    prs.save("Portfolio_Presentation.pptx")
    print("Présentation complète sauvegardée sous : Portfolio_Presentation.pptx")

if __name__ == "__main__":
    create_portfolio_presentation()

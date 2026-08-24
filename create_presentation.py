import pptx 
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Projet CHSA : Agent d'Accueil Hospitalier"
    slide.placeholders[1].text = "Conception, Industrialisation et Déploiement\nPrésenté par Christophe Mpaga - AI Engineer"

    # Slide 2: Contexte
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Contexte et Enjeux"
    content = slide.placeholders[1]
    content.text = ("• Problématique : Surcharge récurrente des services d'urgence.\n"
                    "• Mission : Fluidifier l'orientation des patients.\n"
                    "• Solution : Agent IA décisionnel pour un triage rapide et fiable.")

    # Slide 3: Méthodologie
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Méthodologie : POC vers PROD"
    content = slide.placeholders[1]
    content.text = ("• POC (4 sem) : Faisabilité technique, modèle Qwen3-1.7B, CI/CD.\n"
                    "• MVP (3-6 mois) : Validation clinique, Dashboard MLflow.\n"
                    "• PROD (6-12 mois+) : Haute disponibilité, système multi-agent, intégration SIH (FHIR).")

    # Slide 4: Architecture (Le pivot)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Pivot Architectural"
    content = slide.placeholders[1]
    content.text = ("• Défi : Échec du design monolithique (timeout 240s sur vLLM).\n"
                    "• Pivot : Passage à une architecture modulaire multi-stage.\n"
                    "• Résultat : Découplage API, Modèle (vLLM), et UI sur Cloud Run.\n"
                    "• Bénéfice : Observabilité accrue et maintenance indépendante.")

    # Slide 5: Industrialisation & MLOps
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Industrialisation & MLOps"
    content = slide.placeholders[1]
    content.text = ("• Observabilité : Monitoring de drift via MLflow.\n"
                    "• Sécurité : Intégration de tests de sécurité automatisés (red-teaming).\n"
                    "• Qualité : Logs d'audit système et validation clinique.")

    # Slide 6: Roadmap Agentique
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Évolutions : Architecture Agentique"
    content = slide.placeholders[1]
    content.text = ("• Cerveau : LLM fine-tuné.\n"
                    "• Orchestrateur : LangGraph pour la gestion des flux complexes.\n"
                    "• Outils : API métier, anonymisation, suivi lits/personnels.\n"
                    "• UI : Interface adaptative temps réel.")

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion"
    content = slide.placeholders[1]
    content.text = ("• Bilan : Maîtrise de la chaîne de valeur IA (IAOps).\n"
                    "• Vision : Ingénierie orientée production et conformité.\n"
                    "• Questions ?")

    prs.save("Presentation_CHSA.pptx")
    print("Présentation sauvegardée sous : Presentation_CHSA.pptx")

if __name__ == "__main__":
    create_presentation()
